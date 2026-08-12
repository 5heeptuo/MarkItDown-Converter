"""转换核心测试：MarkItDown 多格式 + OCR fallback + 损坏文件 + 中文路径。

耗时操作（Magika 模型加载、RapidOCR 首次加载）在各测试间共享进程缓存。
"""

from __future__ import annotations

import json
from pathlib import Path

import pytest

from converter import ConvertError, Status, convert_file


# --------------------------------------------------------------------------
# fixtures：生成各类测试文件
# --------------------------------------------------------------------------
def make_docx(path: Path) -> None:
    import docx

    d = docx.Document()
    d.add_heading("测试文档", level=1)
    d.add_paragraph("这是一个用于转换测试的 Word 文档。")
    d.save(path)


def make_xlsx(path: Path) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "数据"
    ws.append(["名称", "数量"])
    ws.append(["苹果", 3])
    ws.append(["香蕉", 5])
    wb.save(path)


def make_pptx(path: Path) -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "演示文稿标题"
    slide.placeholders[1].text = "这是正文内容"
    prs.save(path)


def make_text_image(path: Path, text: str) -> None:
    from PIL import Image, ImageDraw, ImageFont

    img = Image.new("RGB", (900, 220), "white")
    draw = ImageDraw.Draw(img)
    font = None
    for candidate in ("C:/Windows/Fonts/msyh.ttc", "C:/Windows/Fonts/simhei.ttf",
                      "C:/Windows/Fonts/arial.ttf"):
        try:
            font = ImageFont.truetype(candidate, 48)
            break
        except OSError:
            continue
    draw.text((20, 30), text, fill="black", font=font or ImageFont.load_default())
    img.save(path)


def make_scanned_pdf(path: Path) -> None:
    """生成“扫描版”PDF：页面就是一张文字图片，没有文本层。"""
    from PIL import Image
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas

    img_path = path.with_suffix(".probe.png")
    make_text_image(img_path, "扫描讲义 操作系统 Operating Systems")
    try:
        c = canvas.Canvas(str(path))
        c.setPageSize((900, 220))
        c.drawImage(ImageReader(str(img_path)), 0, 0, width=900, height=220)
        c.showPage()
        c.save()
    finally:
        img_path.unlink(missing_ok=True)


# --------------------------------------------------------------------------
# MarkItDown 多格式
# --------------------------------------------------------------------------
@pytest.mark.parametrize(
    ("ext", "maker"),
    [
        (".docx", make_docx),
        (".xlsx", make_xlsx),
        (".pptx", make_pptx),
    ],
)
def test_markitdown_document_formats(tmp_path: Path, ext: str, maker) -> None:
    src = tmp_path / f"文档{ext}"
    dst = tmp_path / f"文档{ext}.md"
    maker(src)
    status, _ = convert_file(src, dst)
    assert status == Status.SUCCESS
    text = dst.read_text(encoding="utf-8")
    assert text.strip(), f"{ext} 转换结果为空"


@pytest.mark.parametrize(
    ("ext", "content"),
    [
        (".txt", "纯文本内容 hello"),
        (".csv", "名称,数量\n苹果,3\n"),
        (".json", json.dumps({"名称": "测试", "值": 42}, ensure_ascii=False)),
        (".html", "<html><body><h1>网页标题</h1><p>网页正文内容</p></body></html>"),
    ],
)
def test_markitdown_text_formats(tmp_path: Path, ext: str, content: str) -> None:
    src = tmp_path / f"文件{ext}"
    dst = tmp_path / f"文件{ext}.md"
    src.write_text(content, encoding="utf-8")
    status, _ = convert_file(src, dst)
    assert status == Status.SUCCESS
    assert dst.read_text(encoding="utf-8").strip()


# --------------------------------------------------------------------------
# 图片 / 扫描 PDF OCR fallback
# --------------------------------------------------------------------------
def test_image_ocr_fallback_chinese(tmp_path: Path) -> None:
    """MarkItDown 对图片输出为空 -> OCR fallback 必须识别出中文。"""
    src = tmp_path / "照片.png"
    dst = tmp_path / "照片.png.md"
    make_text_image(src, "计算机工程 Operating Systems")
    status, _ = convert_file(src, dst)
    assert status == Status.SUCCESS
    text = dst.read_text(encoding="utf-8")
    assert "计算机" in text


def test_image_ocr_fallback_spanish(tmp_path: Path) -> None:
    src = tmp_path / "factura.png"
    dst = tmp_path / "factura.png.md"
    make_text_image(src, "Factura Precio Total")
    status, _ = convert_file(src, dst)
    assert status == Status.SUCCESS
    text = dst.read_text(encoding="utf-8")
    assert "Factura" in text


def test_scanned_pdf_ocr_fallback(tmp_path: Path) -> None:
    """扫描版 PDF（无文本层）必须通过渲染 + OCR 提取文字。"""
    src = tmp_path / "扫描讲义.pdf"
    dst = tmp_path / "扫描讲义.pdf.md"
    make_scanned_pdf(src)
    status, _ = convert_file(src, dst)
    assert status == Status.SUCCESS
    text = dst.read_text(encoding="utf-8")
    assert "操作系统" in text or "Operating" in text


def test_blank_image_reports_no_content(tmp_path: Path) -> None:
    """纯白图片（无文字）-> 不应算成功，应报“无内容”。"""
    from PIL import Image

    src = tmp_path / "空白.png"
    dst = tmp_path / "空白.png.md"
    Image.new("RGB", (400, 200), "white").save(src)
    with pytest.raises(ConvertError) as excinfo:
        convert_file(src, dst)
    assert excinfo.value.category == "no_content"


def test_svg_text_extraction(tmp_path: Path) -> None:
    """SVG 属于 OCR 图片格式，但本质是 XML——应能提取内嵌 <text>。"""
    src = tmp_path / "diagram.svg"
    dst = tmp_path / "diagram.svg.md"
    src.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" width="200" height="100">'
        '<text x="10" y="30">架构示意图</text>'
        '<text x="10" y="60">System Design</text>'
        "</svg>",
        encoding="utf-8",
    )
    status, _ = convert_file(src, dst)
    assert status == Status.SUCCESS
    text = dst.read_text(encoding="utf-8")
    assert "架构示意图" in text
    assert "System Design" in text


# --------------------------------------------------------------------------
# 损坏 / 伪装文件
# --------------------------------------------------------------------------
def test_corrupt_docx_classified_as_corrupt(tmp_path: Path) -> None:
    """扩展名是 .docx 但内容损坏 -> 归类为“文件可能已损坏”。"""
    import os

    src = tmp_path / "broken.docx"
    dst = tmp_path / "broken.docx.md"
    # 以 ZIP 头开头但内容随机：让 DocxConverter 接受扩展名后在解析时失败
    src.write_bytes(b"PK\x03\x04" + os.urandom(128))
    with pytest.raises(ConvertError) as excinfo:
        convert_file(src, dst)
    assert excinfo.value.category == "corrupt"


def test_missing_file_raises_read_error(tmp_path: Path) -> None:
    with pytest.raises(ConvertError) as excinfo:
        convert_file(tmp_path / "不存在.pdf", tmp_path / "不存在.md")
    assert excinfo.value.category == "read"


# --------------------------------------------------------------------------
# 中文 / Unicode 路径
# --------------------------------------------------------------------------
def test_unicode_path_conversion(tmp_path: Path) -> None:
    folder = tmp_path / "大学资料" / "操作系统"
    folder.mkdir(parents=True)
    src = folder / "第一章 笔记.txt"
    dst = folder / "第一章 笔记.md"
    src.write_text("操作系统课程笔记", encoding="utf-8")
    status, _ = convert_file(src, dst)
    assert status == Status.SUCCESS
    assert dst.exists()
    assert "操作系统" in dst.read_text(encoding="utf-8")
